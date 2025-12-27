import os
import shutil
import json
import uuid  
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse  
from pydantic import BaseModel
from gtts import gTTS  
from pptx import Presentation 
from config import UPLOAD_DIR
from rag_engine import rag_system  

app = FastAPI(title="LlamaIndex MongoDB RAG API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class QueryRequest(BaseModel):
    question: str

class QueryResponse(BaseModel):
    answer: str

class TextToSpeechRequest(BaseModel):
    text: str


@app.get("/")
def root():
    return {"message": "start"}

@app.get("/generate/flashcards")
def get_flashcards():
    try:
        json_response = rag_system.generate_flashcards()
        data = json.loads(json_response)
        
        if isinstance(data, dict) and "error" in data:
            return data 

        return {"flashcards": data}

    except Exception as e:
        print(f"Server Error: {e}")
        return {"error": str(e)}
    
@app.get("/generate/quiz")
def get_quiz():
    try:
        json_response = rag_system.generate_quiz()
        data = json.loads(json_response)
        return {"quiz": data}
    except Exception as e:
        return {"error": str(e)}
    
@app.get("/generate/mindmap")
def get_mindmap():
    try:
        json_str = rag_system.generate_interactive_map()
        data = json.loads(json_str)
        return {"interactive_data": data}
    except Exception as e:
        return {"error": str(e)}
    
@app.get("/generate/summary")
def get_summary():
    try:
        summary_text = rag_system.generate_summary()
        return {"summary": summary_text}
    except Exception as e:
        return {"error": str(e)}

@app.post("/generate/audio")
async def generate_audio(request: TextToSpeechRequest):
    try:
        file_name = f"audio_{uuid.uuid4()}.mp3"
        file_path = os.path.join(UPLOAD_DIR, file_name)
        
        tts = gTTS(text=request.text, lang='en', slow=False)
        tts.save(file_path)
        
        return FileResponse(file_path, media_type="audio/mpeg", filename=file_name)
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    
    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Could not save file: {e}")

    try:
        rag_system.add_pdf(file_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PDF: {e}")

    return {"message": "File upload", "filename": file.filename}

@app.post("/ask", response_model=QueryResponse)
def ask_question(request: QueryRequest):
    try:
        response = rag_system.query(request.question)
        return {"answer": response}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error data: {e}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
    
@app.post("/generate/slides")
def generate_slides():
    """API لإنشاء ملف PowerPoint"""
    try:
        json_str = rag_system.generate_slide_content()
        slides_data = json.loads(json_str)
        
        prs = Presentation()
        
        for i, slide_content in enumerate(slides_data):
            if i == 0:
                slide_layout = prs.slide_layouts[0] 
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_content.get("title", "Presentation")
                if slide_content.get("points"):
                    slide.placeholders[1].text = slide_content["points"][0]
            
            else:
                slide_layout = prs.slide_layouts[1] 
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_content.get("title", "Topic")
                
                tf = slide.placeholders[1].text_frame
                for point in slide_content.get("points", []):
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0

        file_name = f"presentation_{uuid.uuid4()}.pptx"
        file_path = os.path.join(UPLOAD_DIR, file_name)
        prs.save(file_path)
        
        return FileResponse(file_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="presentation.pptx")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))